import streamlit as st
import cv2
import numpy as np
import os
import tempfile
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, UnidentifiedImageError

# ---------- Utility functions ----------

def load_image(path, max_size=3000):
    """Load any image type with Pillow and convert to OpenCV format."""
    try:
        img = Image.open(path).convert("RGB")
        # resize if too large
        if img.width > max_size or img.height > max_size:
            img.thumbnail((max_size, max_size))
        return cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
    except UnidentifiedImageError:
        st.warning(f"Skipped (unsupported format): {os.path.basename(path)}")
        return None
    except Exception as e:
        st.warning(f"Skipped {os.path.basename(path)}: {e}")
        return None

def is_tilted(image):
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        lines = cv2.HoughLines(edges, 1, np.pi/180, 200)
        if lines is None:
            return False, 0
        angles = []
        for rho, theta in lines[:,0]:
            angle = (theta * 180 / np.pi) - 90
            if -45 < angle < 45:
                angles.append(angle)
        if len(angles) == 0:
            return False, 0
        median_angle = np.median(angles)
        return abs(median_angle) > 5, median_angle
    except Exception:
        return False, 0

def rotate_image(image, angle):
    (h, w) = image.shape[:2]
    center = (w // 2, h // 2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    return cv2.warpAffine(image, M, (w, h))

def detect_blur_regions(image):
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        lap = cv2.Laplacian(gray, cv2.CV_64F)
        lap_var = cv2.convertScaleAbs(lap)
        _, mask = cv2.threshold(lap_var, 15, 255, cv2.THRESH_BINARY_INV)
        contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        return contours
    except Exception:
        return []

def detect_dark_bright_regions(image):
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        dark_mask = cv2.adaptiveThreshold(gray, 255,
                                          cv2.ADAPTIVE_THRESH_MEAN_C,
                                          cv2.THRESH_BINARY_INV, 15, 10)
        bright_mask = cv2.adaptiveThreshold(gray, 255,
                                            cv2.ADAPTIVE_THRESH_MEAN_C,
                                            cv2.THRESH_BINARY, 15, -10)
        mask = cv2.bitwise_or(dark_mask, bright_mask)
        contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        return contours
    except Exception:
        return []

def analyze_image(image):
    processed = image.copy()
    reasons = []

    # Tilt correction
    tilted, angle = is_tilted(image)
    if tilted:
        image = rotate_image(image, angle)
        processed = image.copy()
        reasons.append(f"Tilted image corrected by {angle:.2f} degrees")

    # Blurry regions
    blur_contours = detect_blur_regions(image)
    if blur_contours:
        for c in blur_contours:
            if cv2.contourArea(c) > 500:
                x, y, w, h = cv2.boundingRect(c)
                cv2.rectangle(processed, (x, y), (x+w, y+h), (0, 0, 255), 2)
        reasons.append("Blurry regions detected")

    # Dark/Bright regions
    light_contours = detect_dark_bright_regions(image)
    if light_contours:
        for c in light_contours:
            if cv2.contourArea(c) > 500:
                x, y, w, h = cv2.boundingRect(c)
                cv2.rectangle(processed, (x, y), (x+w, y+h), (255, 0, 0), 2)
        reasons.append("Dark or overexposed regions detected")

    if not reasons:
        reasons.append("No significant edge case detected")

    return processed, reasons

def create_ppt(results, output_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for orig_path, orig_img, proc_img, reasons in results:
        try:
            slide = prs.slides.add_slide(blank_slide_layout)

            orig_temp = orig_path + "_orig_tmp.jpg"
            proc_temp = orig_path + "_proc_tmp.jpg"
            cv2.imwrite(orig_temp, orig_img)
            cv2.imwrite(proc_temp, proc_img)

            left = Inches(0.5)
            top = Inches(1)
            slide.shapes.add_picture(orig_temp, left, top, width=Inches(4.5))

            left = Inches(5.5)
            slide.shapes.add_picture(proc_temp, left, top, width=Inches(4.5))

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
            tf = txBox.text_frame
            for reason in reasons:
                p = tf.add_paragraph()
                p.text = reason
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(255, 0, 0)

            os.remove(orig_temp)
            os.remove(proc_temp)
        except Exception as e:
            st.warning(f"Could not add {os.path.basename(orig_path)} to PPT: {e}")

    prs.save(output_path)

# ---------- Streamlit App ----------

st.title("Edge Case Detection and PPT Generator")

uploaded_zip = st.file_uploader("Upload a ZIP containing images", type=["zip"])

if uploaded_zip:
    with tempfile.TemporaryDirectory() as tmpdir:
        zip_path = os.path.join(tmpdir, "uploaded.zip")
        with open(zip_path, "wb") as f:
            f.write(uploaded_zip.read())

        with zipfile.ZipFile(zip_path, "r") as zip_ref:
            zip_ref.extractall(tmpdir)

        results = []
        for root, _, files in os.walk(tmpdir):
            for file in files:
                path = os.path.join(root, file)
                img = load_image(path)
                if img is None:
                    continue
                try:
                    processed, reasons = analyze_image(img)
                    results.append((path, img, processed, reasons))
                except Exception as e:
                    st.warning(f"Error analyzing {os.path.basename(path)}: {e}")

        if results:
            output_ppt = os.path.join(tmpdir, "edge_cases.pptx")
            create_ppt(results, output_ppt)
            st.success("PPT created successfully!")
            with open(output_ppt, "rb") as f:
                st.download_button("Download PPT", f, file_name="edge_cases.pptx")
        else:
            st.warning("No valid images processed from the uploaded ZIP.")
